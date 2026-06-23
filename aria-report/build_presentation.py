#!/usr/bin/env python3
"""Build a PowerPoint presentation from presentation-restructured.md.

Requires: python-pptx (install: pip install python-pptx)
Run from the aria-report directory:
    python3 build_presentation.py
"""
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

MD_FILE = "presentation-restructured.md"
OUTPUT = "presentation-restructured.pptx"

# Layout dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Theme colors
COLOR_TITLE = RGBColor(0x1A, 0x23, 0x7E)  # dark blue
COLOR_TEXT = RGBColor(0x22, 0x22, 0x22)
COLOR_ACCENT = RGBColor(0x00, 0x6D, 0xA8)


def split_slides(md_path):
    with open(md_path, "r", encoding="utf-8") as f:
        text = f.read()
    # Drop YAML front matter if present
    if text.startswith("---"):
        parts = text.split("---", 2)
        if len(parts) >= 3:
            text = parts[2]
    raw_slides = re.split(r"\n---\s*\n", text.strip())
    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        lines = raw.splitlines()
        # First non-empty line should be # Slide N — Title
        title = "Slide"
        body = []
        for i, line in enumerate(lines):
            s = line.strip()
            if i == 0 and s.startswith("# "):
                title = s[2:].strip()
                continue
            body.append(line)
        slides.append((title, "\n".join(body).strip()))
    return slides


def add_title_slide(prs, title_text):
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    # Background bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(3.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_TITLE
    shape.line.fill.background()
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12.1), Inches(1.4))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return slide


def add_content_slide(prs, title_text, body):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_TITLE
    bar.line.fill.background()
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Parse body for image first
    image_match = re.search(r"!\[([^\]]*)\]\(([^)]+)\)", body)
    if image_match:
        img_path = image_match.group(2)
        if os.path.exists(img_path):
            # Place image on right, text on left
            try:
                slide.shapes.add_picture(img_path, Inches(6.6), Inches(1.4), width=Inches(6.2))
            except Exception as e:
                print(f"  Warning: could not add image {img_path}: {e}")
            text_left = Inches(0.5)
            text_width = Inches(5.9)
            text_top = Inches(1.3)
            text_height = Inches(5.9)
        else:
            print(f"  Warning: image not found {img_path}")
            text_left = Inches(0.5)
            text_width = Inches(12.3)
            text_top = Inches(1.3)
            text_height = Inches(5.9)
        # Remove image markdown line from body
        body = re.sub(r"!\[([^\]]*)\]\(([^)]+)\)\n?", "", body).strip()
    else:
        text_left = Inches(0.5)
        text_width = Inches(12.3)
        text_top = Inches(1.3)
        text_height = Inches(5.9)

    # Parse body into paragraphs / tables / code
    add_body_to_slide(slide, body, text_left, text_top, text_width, text_height)
    return slide


def add_body_to_slide(slide, body, left, top, width, height):
    lines = body.splitlines()
    # If first non-empty line is a heading, render it as subtitle
    if lines and lines[0].strip().startswith("#"):
        subtitle = lines[0].strip().lstrip("#").strip()
        sub_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        top += Inches(0.55)
        height -= Inches(0.55)
        lines = lines[1:]

    # Detect and render markdown table
    table_match = None
    for i, line in enumerate(lines):
        if "|" in line and i + 1 < len(lines) and re.match(r"^\|?[-:\s|]+-?\|?\s*$", lines[i + 1]):
            table_match = extract_table(lines[i:])
            break

    if table_match:
        headers, rows, consumed = table_match
        # Render table
        table_height = min(Inches(0.6 + 0.45 * (len(rows) + 1)), height)
        table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, table_height).table
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_ACCENT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, val in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = COLOR_TEXT
        top += table_height + Inches(0.2)
        height -= table_height + Inches(0.2)
        # Remove table lines from body
        lines = lines[:i] + lines[i + consumed:]

    # Detect code block
    code_blocks = list(re.finditer(r"```(?:\w+)?\n(.*?)\n```", body, re.DOTALL))
    if code_blocks:
        # We'll render remaining lines as bullets, code blocks separately below
        pass

    # Render remaining text as bullet list / paragraphs
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    for line in lines:
        s = line.rstrip()
        if not s.strip():
            continue
        if s.strip().startswith("```"):
            continue
        if re.match(r"^\|?[-:\s|]+-?\|?\s*$", s):
            continue
        p = tf.add_paragraph()
        # Heading inside body
        if s.startswith("## "):
            p.text = s[3:].strip()
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = COLOR_ACCENT
        elif s.startswith("### "):
            p.text = s[4:].strip()
            p.font.size = Pt(16)
            p.font.bold = True
        elif s.startswith("- ") or s.startswith("* "):
            p.text = "• " + s[2:].strip()
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_TEXT
        elif re.match(r"^\s+- ", s):
            p.text = "◦ " + s.strip()[2:].strip()
            p.level = 1
            p.font.size = Pt(14)
        elif re.match(r"^\d+\.\s", s):
            p.text = s.strip()
            p.font.size = Pt(15)
        elif s.strip().startswith(">"):
            p.text = s.strip().lstrip(">").strip()
            p.font.italic = True
            p.font.size = Pt(15)
            p.font.color.rgb = COLOR_ACCENT
        else:
            p.text = s.strip()
            p.font.size = Pt(15)
            p.font.color.rgb = COLOR_TEXT

    # Render code blocks as text boxes below main text if space allows
    for m in code_blocks:
        code = m.group(1)
        code_box = slide.shapes.add_textbox(left, Inches(5.0), width, Inches(2.2))
        ctf = code_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = code
        cp.font.name = "Consolas"
        cp.font.size = Pt(9)
        cp.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def extract_table(lines):
    """Extract markdown table from lines. Return (headers, rows, consumed_count)."""
    headers = [c.strip() for c in lines[0].split("|") if c.strip()]
    rows = []
    consumed = 2  # header + separator
    for line in lines[2:]:
        if "|" not in line:
            break
        cells = [c.strip() for c in line.split("|")]
        # Drop empty leading/trailing cells
        cells = [c for c in cells if c or c == ""]
        if not any(cells):
            break
        rows.append(cells[: len(headers)])
        consumed += 1
    return headers, rows, consumed


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slides = split_slides(MD_FILE)
    # Drop trailing empty / marker slides
    slides = [(t, b) for t, b in slides if not (t == "Slide" and (not b or "End of Presentation" in b))]
    print(f"Found {len(slides)} slides in {MD_FILE}")

    for idx, (title, body) in enumerate(slides, start=1):
        # Extract just the title text after "Slide N — "
        m = re.match(r"Slide\s+\d+\s*[-–—]\s*(.+)", title)
        title_text = m.group(1).strip() if m else title
        print(f"Building slide {idx}: {title_text}")

        if idx == 1 and "Title" in title:
            slide = add_title_slide(prs, title_text)
            # Add subtitle and author info from body on title slide
            text_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(12.1), Inches(3.5))
            tf = text_box.text_frame
            tf.word_wrap = True
            for line in body.splitlines():
                s = line.strip()
                if not s:
                    continue
                p = tf.add_paragraph()
                p.text = s.lstrip("#").strip()
                p.font.size = Pt(20)
                p.font.color.rgb = COLOR_TEXT
                p.alignment = PP_ALIGN.CENTER
        else:
            add_content_slide(prs, title_text, body)

    prs.save(OUTPUT)
    print(f"Saved {OUTPUT}")


if __name__ == "__main__":
    main()
